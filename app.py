import streamlit as st
import openai
import re
import textwrap
from io import BytesIO
from docx import Document
from fpdf import FPDF

import random

st.set_page_config(page_title="Super Teacher", layout="wide")


import random
import re

def create_cloze(passage: str, num_blanks: int = 5):
    stopwords = {"the", "a", "an", "and", "or", "in", "on", "of", "to", "for", "with", "from", "at", "by", "as", "that"}
    
    words = re.findall(r'\b\w+\b', passage)
    candidates = [w for w in set(words) if w.lower() not in stopwords and len(w) > 3]

    if len(candidates) < num_blanks:
        num_blanks = len(candidates)

    selected = random.sample(candidates, num_blanks)
    answer_map = {word: "_____" for word in selected}


    def replacer(match):
        word = match.group(0)
        if word in answer_map:
            replacement = answer_map.pop(word)
            return replacement
        return word

    pattern = re.compile(r'\b(' + '|'.join(re.escape(w) for w in selected) + r')\b')
    cloze_passage = pattern.sub(replacer, passage)
    random.shuffle(selected)
    return cloze_passage, selected



# Initialize the client using the older pattern
import os  # if not already at top
client = openai.OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
if not client.api_key:
    st.warning("Please enter your OpenAI API key in the Streamlit secrets.")
    st.stop()

# --- SIDEBAR: TOOL SELECTION ---
st.sidebar.title("SIDEKICK")
st.sidebar.title("Kicking goals like a champ✨")
tool = st.sidebar.radio(
    "Choose a tool:",
    ["Lesson Builder", "Unit Planner", "Unit Glossary Generator", "Worksheet Generator", "Test Creator", "Email Assistant", "Video Assistant", "Feedback Assistant", "Self Care Tool", "Feeling Peckish"]
)


# ========== HELPER FUNCTION ==========
def chat_completion_request(system_msg, user_msg, max_tokens=1000, temperature=0.7):
    """
    A helper to call GPT-3.5-turbo with system & user messages using the client.
    """
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_msg}
        ],
        max_tokens=max_tokens,
        temperature=temperature
    )
    return response.choices[0].message.content.strip()

import re

def display_output_block(text):
    cleaned_text = text.strip().replace("*", "").replace("#", "")
    html_text = cleaned_text.replace("\n", "<br>")  # Move this out of the f-string
    st.markdown(
        f"""
        <div style='background-color: white; color: black; padding: 20px; 
                    border-radius: 8px; font-family: sans-serif; 
                    font-size: 16px; line-height: 1.6; white-space: pre-wrap; margin-left: 0;'>
            {html_text}
        </div>
        """,
        unsafe_allow_html=True
    )







# ========== TOOL 1: LESSON BUILDER ==========
if tool == "Lesson Builder":
    st.markdown(
    """
    <h1 style='font-size: 2.25rem; font-weight: 600; margin-bottom: 1rem;
               font-family: "Segoe UI", sans-serif;'>📝 Lesson Builder</h1>
    """,
    unsafe_allow_html=True
)

    year = st.text_input("Grade Level (e.g. 7)", placeholder="Enter grade level here")
    subject = st.text_input("Subject (e.g. English, Science)")
    topic = st.text_input("Lesson Topic")
    duration = st.slider("Lesson Duration (minutes)", 30, 120, 70, step=5)
    lesson_count = st.number_input("Number of Lessons", min_value=1, value=1, step=1)
    goal_focus = st.selectbox("Lesson Focus", ["Skills-Based", "Knowledge-Based", "Critical Thinking", "Creative Thinking"])
    include_curriculum = st.checkbox("Include V9 curriculum reference")
    device_use = st.multiselect("Resources Available to Students", ["Laptop/Tablet", "Textbook", "Worksheet"])
    grouping = st.selectbox("Grouping Preference", ["Individual", "Pairs", "Small Groups", "Whole Class"])
    lesson_style = st.selectbox("Lesson Style", ["Quiet/Reflective", "Discussion-Based", "Hands on", "Creative"])
    assessment = st.selectbox("Assessment Format", ["No Assessment", "Exit Slip", "Short Response", "Group Presentation", "Quiz"])
    differentiation = st.multiselect("Include Differentiation for:", ["Support", "Extension", "ESL", "Neurodiverse"])

    # After your input fields are defined
    if st.button("Generate Lesson Plan"):
        prompt_parts = [
            f"Create {lesson_count} lesson(s), each {duration} minutes long, for a Year {year} {subject} class on '{topic}'.",
            f"Start each lesson with a clear Learning Goal aligned to a {goal_focus.lower()} outcome.",
            "Structure each lesson with: Hook, Learning Intentions, Warm-up, Main Task, Exit Ticket.",
            f"The lesson should use {', '.join(device_use).lower()}. Students should work in {grouping.lower()}.",
            f"Use a {lesson_style.lower()} approach."
        ]
      
        if differentiation:
            prompt_parts.append("Include differentiation strategies for: " + ", ".join(differentiation) + ".")
        if assessment != "No Assessment":
            prompt_parts.append(f"End each lesson with a {assessment.lower()} as an assessment.")
        if include_curriculum:
            prompt_parts.append("Align the lesson with the Australian V9 curriculum.")

        full_prompt = " ".join(prompt_parts)

        with st.spinner("Planning your lesson(s)..."):
            lesson_plan = chat_completion_request(
                system_msg="You are a practical, creative Australian teacher.",
                user_msg=full_prompt,
                max_tokens=1200
            )
            formatted_plan = lesson_plan.replace("* ", "• ")
            formatted_plan = re.sub(r"^#+\s*(.+)$", r"<br><b>\1</b>", formatted_plan, flags=re.MULTILINE)
            html_plan = formatted_plan.replace("\n", "<br>")
            st.markdown(
                f"""
                <div style='background-color: #f9f9f9; padding: 20px; border-radius: 8px;
                            font-family: sans-serif; font-size: 16px; color: #111;
                            line-height: 1.6; white-space: pre-wrap;'>
                    {html_plan}
                </div>
                """,
                unsafe_allow_html=True
            )


        # After displaying the lesson plan:
        st.subheader("Export Options")
        
        # New: Download as PowerPoint Slides (Multiple Slides Version)
        from pptx import Presentation
        from pptx.util import Inches, Pt

        # Clean the lesson_plan to remove asterisks (*) and hashes (#)
        export_plan = re.sub(r'[\*\#]', '', lesson_plan)
        
        ppt_buffer = BytesIO()
        prs = Presentation()
        
        # Split the lesson plan text into chunks. For example, splitting by double newlines:
        slide_sections = re.split(r'\n\s*\n', lesson_plan.strip())
        
        for i, section in enumerate(slide_sections, start=1):
            # Select a blank slide layout (adjust as desired)
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            
            # Optional: add a title for each slide
            title_placeholder = slide.shapes.title
            if title_placeholder:
                title_placeholder.text = f"Slide {i}"
            
            # Add a text box containing the section text
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            tf = text_box.text_frame
            tf.text = section  # Put the section text here
            tf.word_wrap = True # Enable automatic wrapping
            
            # Format the text (if desired)
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
        
        # Save the presentation to the buffer and prepare for download
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        st.download_button(
            label="📊 Download PowerPoint",
            data=ppt_buffer,
            file_name="lesson_plan.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )


# ========== TOOL 2: FEEDBACK ASSISTANT ==========
elif tool == "Feedback Assistant":
    st.header("🧠 Feedback Assistant")
    student_text = st.text_area("Paste student writing here:")
    tone = st.selectbox("Choose feedback tone", ["Gentle", "Firm", "Colloquial"])

    if st.button("Generate Feedback"):
        feedback_prompt = (
            f"Give feedback on the following student writing using the Star and 2 Wishes model. "
            f"Highlight errors in spelling, grammar, cohesion, repetition, and sentence structure by surrounding them with **bold**. "
            f"Use a {tone.lower()} tone.\n\n"
            f"Student text:\n{student_text}"
        )
        with st.spinner("Analysing writing..."):
            feedback = chat_completion_request(
                system_msg="You are a kind, helpful teacher giving writing feedback.",
                user_msg=feedback_prompt,
                max_tokens=800
            )
           
            display_output_block(feedback)


# ========== TOOL 3: EMAIL ASSISTANT ==========
elif tool == "Email Assistant":
    st.header("✉️ Email Assistant")
    recipient = st.selectbox("Who is the email to?", ["Parent", "Student", "Staff", "Other"])
    context = st.text_area("Briefly describe the situation or what you want to say:")
    tone = st.selectbox("Choose tone", ["Supportive", "Professional", "Friendly"])

    if st.button("Generate Email"):
        email_prompt = (
            f"Write a {tone.lower()} email to a {recipient.lower()} about the following situation:\n"
            f"{context}\n\n"
            f"Keep the email clear, respectful, and well-structured."
        )
        with st.spinner("Writing your email..."):
            email_content = chat_completion_request(
                system_msg="You are an experienced teacher writing professional school emails.",
                user_msg=email_prompt,
                max_tokens=800
            )
            
            display_output_block(email_content)

# ========== TOOL 4: UNIT GLOSSARY GENERATOR ==========
elif tool == "Unit Glossary Generator":
    st.header("📘 Unit Glossary Generator")
    year = st.text_input("Grade Level (e.g. 7)", placeholder="Enter grade level here")
    subject = st.text_input("Subject (e.g. Science, HASS)")
    topic = st.text_input("Topic or Unit Focus (e.g. Body Systems, Volcanoes)")
   

    if st.button("Generate Glossary"):
        glossary_prompt = (
            f"Create a 3-tier vocabulary glossary for a Year {year} {subject} unit on '{topic}'. "
            "Use this structure:\n"
            "Tier 1 (General): 10 basic words students must know.\n"
            "Tier 2 (Core): 7 subject-specific words they will encounter in lessons.\n"
            "Tier 3 (Stretch): 5 challenge words that extend thinking.\n\n"
            "Use bullet points. Keep each definition under 20 words. "
            "Use student-friendly language, especially for primary year levels."
        )
        with st.spinner("Generating vocabulary list..."):
            glossary = chat_completion_request(
                system_msg="You are a helpful and experienced curriculum-aligned teacher.",
                user_msg=glossary_prompt,
                max_tokens=700
            )
            
            display_output_block(glossary)

# ========== TOOL 5: UNIT PLANNER ==========
elif tool == "Unit Planner":
    st.header("📘 Unit Planner")
    year = st.text_input("Grade Level (e.g. 7)", placeholder="Enter grade level here")
    subject = st.text_input("Subject (e.g. HASS, English, Science)")
    topic = st.text_input("Unit Topic or Focus (e.g. Ancient Egypt, Persuasive Writing)")
    weeks = st.slider("Estimated Duration (Weeks)", 1, 10, 5)

    include_assessment = st.checkbox("Include Assessment Suggestions?")
    include_hook = st.checkbox("Include Hook Ideas for Lesson 1?")
    include_fast_finishers = st.checkbox("Include Fast Finisher Suggestions?")
    include_cheat_sheet = st.checkbox("Include Quick Content Cheat Sheet (for teacher)?")

    # Use session_state to store the generated plan so it doesn't reset on download clicks
    if "unit_plan" not in st.session_state:
        st.session_state["unit_plan"] = None

    if st.button("Generate Unit Plan"):
        prompt_parts = [
            f"Create a unit plan overview for a Year {year} {subject} unit on '{topic}'.",
            f"The unit runs for approximately {weeks} weeks.",
            "Include the following sections:",
            "1. A short Unit Overview (what it's about).",
            "2. 3–5 clear Learning Intentions.",
            "3. A suggested sequence of subtopics or concepts to explore each week.",
            "4. A comprehensive list of lesson types or activity ideas that would suit this unit."
        ]
        if include_assessment:
            prompt_parts.append("5. Include 1–2 assessment ideas (format only, keep it brief).")
        if include_hook:
            prompt_parts.append("6. Suggest 2–3 engaging Hook Ideas for Lesson 1.")
        if include_fast_finishers:
            prompt_parts.append("7. Suggest Fast Finisher or Extension Task ideas.")
        if include_cheat_sheet:
            prompt_parts.append("8. Provide a Quick Content Cheat Sheet: 10 bullet-point facts a teacher should know to teach this unit.")

        full_prompt = " ".join(prompt_parts)

        with st.spinner("Planning your unit..."):
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a practical and experienced curriculum-aligned teacher in Australia."},
                    {"role": "user", "content": full_prompt}
                ]
            )

        if response and response.choices:
            st.session_state["unit_plan"] = response.choices[0].message.content
        else:
            st.warning("⚠️ Unit plan generation failed. Please try again.")

    # If the unit plan is generated, show it and provide download options
    if st.session_state["unit_plan"]:
        unit_plan = st.session_state["unit_plan"]
        st.subheader("Your Generated Unit Plan")
        display_output_block(unit_plan)
        st.markdown("---")
        st.subheader("📄 Export Options")

        
        # Remove all asterisks and hashes for exports
        export_plan = re.sub(r'[\*\#]', '', unit_plan)


         # ========== WORD EXPORT ==========
        word_buffer = BytesIO()
        doc = Document()
        doc.add_paragraph(export_plan)
        # Remove document protection if it exists to avoid a locked/read-only file
        try:
            protection = doc.settings.element.xpath('//w:documentProtection')
            if protection:
                protection[0].getparent().remove(protection[0])
        except Exception:
            pass
        doc.save(word_buffer)
        word_buffer.seek(0)
        st.download_button(
            label="📝 Download Word",
            data=word_buffer,
            file_name="unit_plan.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            key="word_download_btn"
        )


# ========== TOOL 6: WORKSHEET GENERATOR ==========
elif tool == "Worksheet Generator":
    st.header("📝 Worksheet Generator")

    # Input fields
    year = st.text_input("Grade Level (e.g. 7)", placeholder="Enter grade level here")
    learning_goal = st.text_area("Enter a learning goal or paste a lesson plan excerpt", height=200)
    num_questions = st.slider("Number of questions", min_value=3, max_value=15, value=5, step=1)
    passage_length = st.slider("Desired word count for the information passage (50-200)", min_value=50, max_value=200, value=100, step=10)

    # Toggle cloze activity
    cloze_activity = st.checkbox("Make the passage a cloze activity (fill-in-the-blank worksheet)")
    if cloze_activity:
        num_blanks = st.slider("Number of words to remove", min_value=5, max_value=20, value=10, step=1)

    if st.button("Generate Worksheet"):
        if cloze_activity:
            # Step 1: Generate base passage and questions from GPT
            worksheet_prompt = (
                f"Based on the following learning goal or lesson plan excerpt for Year {year}:\n\n"
                f"{learning_goal}\n\n"
                f"Write an information passage of about {passage_length} words. "
                f"Then generate {num_questions} short answer questions for students based on the passage. "
                f"Do not remove any words or create blanks. Do not include headers (eg 'Short Answer Questions'). After the questions, include an 'Answer Key:' section with the correct answers."
            )
        
            with st.spinner("Generating cloze worksheet..."):
                response = chat_completion_request(
                    system_msg="You are a creative teacher assistant who specializes in generating educational worksheets.",
                    user_msg=worksheet_prompt,
                    max_tokens=1000,
                    temperature=0.7
                )
                
            
    
            if "1." in response:
                split_index = response.find("1.")
                passage = response[:split_index].strip()
                questions = response[split_index:].strip()
            else:
                passage = response.strip()
                questions = ""

        
            # Extract body only from the passage
            header_1 = "Information Passage:"
            header_2 = "Short Answer Questions:"
            body_start = passage.find(header_1) + len(header_1)
            body_end = passage.find(header_2)
            if body_end == -1:
                body_end = len(passage)
            body_only = passage[body_start:body_end].strip()
        
            # Create cloze version of passage
            cloze_body, answer_list = create_cloze(body_only, num_blanks=num_blanks)
            cloze_passage = f"{header_1}\n\n{cloze_body}".strip()

        
           # Try to split out GPT answer section
            if "Answer Key:" in questions:
                question_part, answer_part = questions.split("Answer Key:", 1)

                
                # Remove rogue headings from both sections
              
                question_part = re.sub(
                    r"(?im)^.*short\s*answer\s*questions.*\n?",
                    "",
                    question_part
                ).strip()
                
                answer_part = re.sub(
                    r"(?im)^.*short\s*answer\s*questions.*\n?",
                    "",
                    answer_part
                ).strip()
        

            
                questions_only = question_part
                answers_only = answer_part
            else:
                questions_only = questions
                answers_only = ""

        
            # Build cloze answer list
            cloze_answers = "\n".join([f"{i+1}. {word}" for i, word in enumerate(answer_list)])
        
            # Build final worksheet
            worksheet = (
                f"**Cloze Passage:**\n\n{cloze_passage}\n\n"
                f"**Answer Key (Blanks):**\n\n{cloze_answers}\n\n"
                f"{questions_only}"
            )
        
            if answers_only:
                worksheet += f"\n\n**Short Answer Answers:**\n\n{answers_only}"
            # Remove any occurrence of the "Short Answer Questions:" header from the entire worksheet
            worksheet = re.sub(
                r"(?im)^\s*short\s*answer\s*questions\s*[:\-]*\s*\n?", 
                "", 
                worksheet
            ).strip()

        
            display_output_block(worksheet)
            


        else:
            # Regular worksheet
            worksheet_prompt = (
                f"Based on the following learning goal or lesson plan excerpt for Year {year}:\n\n"
                f"{learning_goal}\n\n"
                f"Generate a worksheet containing {num_questions} short answer questions for students. "
                f"The accompanying information passage should be approximately {passage_length} words. "
                f"List all the questions first, then at the bottom provide the corresponding answers."
            )

            with st.spinner("Generating worksheet..."):
                worksheet = chat_completion_request(
                    system_msg="You are a creative teacher assistant who specializes in generating educational worksheets.",
                    user_msg=worksheet_prompt,
                    max_tokens=1000,
                    temperature=0.7
                )
                display_output_block(worksheet)


        # ---- Export Options ----
        st.subheader("Export Options")
        st.write("Do you think your students will notice the answers at the bottom? :)")
        export_worksheet = re.sub(r'[\*\#]', '', worksheet)


        word_buffer = BytesIO()
        doc = Document()
        doc.add_paragraph(export_worksheet)
        try:
            protection = doc.settings.element.xpath('//w:documentProtection')
            if protection:
                protection[0].getparent().remove(protection[0])
        except Exception:
            pass
        doc.save(word_buffer)
        word_buffer.seek(0)
        st.download_button(
            label="📝 Download Word",
            data=word_buffer,
            file_name="worksheet.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
      


# ========== PECKISH ==========
elif tool == "Feeling Peckish":
    st.header("🍽️ Feeling Peckish")
    dish = st.text_input("Enter a dish or beverage (e.g., chicken curry or espresso martini):")
    if st.button("Get Recipe"):
        recipe_prompt = (
            f"Provide a detailed recipe for {dish}. Include a list of ingredients, step-by-step instructions, "
            "and any useful tips for preparation."
        )
        with st.spinner("Messing up the kitchen..."):
            recipe = chat_completion_request(
                system_msg="You are an expert chef who provides creative and detailed recipes.",
                user_msg=recipe_prompt,
                max_tokens=800,
                temperature=0.7
            )
        display_output_block(recipe)

# ========== SELF CARE ==========
elif tool == "Self Care Tool":
    st.header("💖 Self Care Tool")
    # Optionally ask how the user is feeling
    mood = st.selectbox("How are you feeling today?", ["Stressed", "Happy", "Tired", "Lonely", "Motivated", "Calm"])
    if st.button("Get Self Care Tip"):
        self_care_prompt = (
            f"Provide a self care tip for someone who is feeling {mood}. "
            "Include some amusing or uplifting advice about what they can do today, "
            "and tell them something amazing about themselves."
        )
        with st.spinner("Breathe in the zen..."):
            tip = chat_completion_request(
                system_msg="You are a caring self care advisor who offers thoughtful, humorous, and uplifting self care tips.",
                user_msg=self_care_prompt,
                max_tokens=300,
                temperature=0.8
            )
        display_output_block(tip)

# ========== Video Assistant ==========
elif tool == "Video Assistant":
    st.header("🎥 Video Assistant")
    grade = st.text_input("Grade Level (e.g. 7)", placeholder="Enter grade level here")
    video_description = st.text_area("What is the video about? Provide a brief overview:")

    
    if st.button("Generate Content"):
        video_prompt = (
            f"For a Grade {grade} class, generate the following based on the video description:\n"
            f"1. A few discussion starter questions\n"
            f"2. A list of key vocabulary words they might encounter along with a brief definition of each but dont repeat the word in the definition\n"
            f"3. Some thoughtful questions to promote deeper engagement\n\n"
            f"Video description: {video_description}"
        )
        
        with st.spinner("Generating video assistant content..."):
            video_content = chat_completion_request(
                system_msg="You are a creative educational content generator.",
                user_msg=video_prompt,
                max_tokens=500,
                temperature=0.7
            )
            display_output_block(video_content)


# ========== TOOL 7: TEST CREATOR ==========
elif tool == "Test Creator":
    st.header("🧪 Test Creator")

    # Input fields
    year = st.text_input("Grade Level (e.g. 7)", placeholder="Enter grade level here")
    subject = st.text_input("Subject", placeholder="e.g. English, Science, HASS")
    topic = st.text_input("Topic", placeholder="e.g. Fractions, Ancient Rome, Persuasive Texts")


    
    num_tf = st.number_input("Number of True/False Questions (Max 20)", min_value=0, max_value=20, value=3, step=1)
    num_mcq = st.number_input("Number of Multiple Choice Questions (Max 20)", min_value=0, max_value=20, value=3, step=1)
    num_sa = st.number_input("Number of Short Response Questions (Max 5)", min_value=0, max_value=5, value=4, step=1)
    num_er = st.number_input("Number of Extended Response Questions (Max 2)", min_value=0, max_value=2, value=0, step=1)

    
    # Total for prompt
    total_qs = num_tf + num_mcq + num_sa + num_er




    mix_difficulty = st.checkbox("Mix difficulty levels?", value=True)
    include_instructions = st.checkbox("Include instructions at the top?", value=True)
    include_answers = st.checkbox("Generate an answer sheet?", value=True)

    if st.button("Generate Test"):
        test_prompt = (
            f"Create a {total_qs}-question test for Year {year} students on the topic '{topic}' in the subject '{subject}'.\n\n"
        )
        if mix_difficulty:
            test_prompt += "Mix easy, medium, and hard questions.\n"
        if include_instructions:
            test_prompt += "Include clear test instructions at the top.\n"
        if include_answers:
            test_prompt += "After the test, include an 'Answer Sheet:' section with correct answers.\n"

        
        test_prompt = (
            f"Create a test with the following structure for Year {year} students on the topic '{topic}' in {subject}:\n"
            f"- {num_tf} True/False questions\n"
            f"- {num_mcq} Multiple Choice questions\n"
            f"- {num_sa} Short Answer questions (3–5 sentence responses)\n"
            f"- {num_er} Extended Response questions (at least 10 sentences)\n\n"
            "IMPORTANT: You must include ALL of these sections, even if the number of questions is low.\n"
            "Group the questions by type, in the order listed above. Number the questions sequentially (do not reset numbering).\n"
        )
        
        # Add section instructions
        if num_sa > 0:
            test_prompt += "Before the short answer section, include this line: 'Short Answer responses require 3–5 sentences.'\n"
        if num_er > 0:
            test_prompt += "Before the extended response section, include this line: 'Extended Response answers require a well-developed paragraph of at least 10 sentences.'\n"
        
        # Final formatting instructions
        test_prompt += (
            "Add 'Student Name:___________________' at the very top.\n"
            "Use clean formatting suitable for copying into Word. Avoid markdown symbols like asterisks or hashes.\n"
            "Leave space after each question for student responses.\n"
            "End the test with the phrase: 'End of Test'."
        )



        with st.spinner("Generating test..."):
            test_output = chat_completion_request(
                system_msg="You are an expert teacher creating clear, printable classroom tests.",
                user_msg=test_prompt,
                max_tokens=1200,
                temperature=0.7
            )

            display_output_block(test_output)

        # ---- Export Options ----
        st.subheader("Export Options")
        st.write("How many students studied? ;)")

        export_test = re.sub(r'[\*\#]', '', test_output)

        word_buffer = BytesIO()
        doc = Document()
        doc.add_paragraph(export_test)
        try:
            protection = doc.settings.element.xpath('//w:documentProtection')
            if protection:
                protection[0].getparent().remove(protection[0])
        except Exception:
            pass
        doc.save(word_buffer)
        word_buffer.seek(0)
        st.download_button(
            label="📥 Download Word",
            data=word_buffer,
            file_name="test.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )




    

 # Generate a unique Teacher Boost dynamically using ChatGPT (no pre-populated list)
st.sidebar.markdown("<br><hr><br>", unsafe_allow_html=True)  # extra spacing and a divider

# Check if a teacher boost has already been generated for this session
if "teacher_boost" not in st.session_state:
    with st.spinner(""):
        boost_prompt = (
            "Generate an uplifting teacher boost message that is either funny, sarcastic, "
            "or a random quirky fact. Please ensure your answer is 40 tokens or less. "
            "Don't use words like hell or other possible offensive terms. You must not use hell or other offensive words."
        )
        st.session_state["teacher_boost"] = chat_completion_request(
            system_msg="You are a creative teacher boost generator.",
            user_msg=boost_prompt,
            max_tokens=40,
            temperature=0.9
        )

st.sidebar.markdown(f"_{st.session_state['teacher_boost']}_")








